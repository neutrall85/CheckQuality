import os
import logging
from pptx import Presentation

logger = logging.getLogger(__name__)

class PowerpointTemplateManager:
    """Загружает тему из шаблона PowerPoint (.potx или .pptx) без слайдов."""

    def __init__(self, builtin_template_path: str = None):
        self.builtin_template_path = builtin_template_path

    def get_presentation(self, template_path: str = None) -> Presentation:
        path = template_path or self.builtin_template_path
        if path and os.path.exists(path):
            return Presentation(path)
        return Presentation()

    def get_presentation_with_theme(self, template_path: str = None) -> Presentation:
        """Загружает презентацию из шаблона и очищает её от слайдов, сохраняя тему."""
        prs = self.get_presentation(template_path)
        # Безопасное удаление слайдов с обёрткой исключений
        while len(prs.slides) > 0:
            try:
                rId = prs.slides._sldIdLst[0].get('r:id')
                prs.part.drop_rel(rId)
                del prs.slides._sldIdLst[0]
            except Exception as e:
                logger.warning(f'Не удалось удалить слайд из шаблона: {e}')
                break
        return prs