import os
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

from core.models.statistics import Statistics
from core.models.config import ConfigModel
from .base_report import BaseReport
from .powerpoint_template_manager import PowerpointTemplateManager
from .placeholder_renderer import PlaceholderRenderer
from .chart_builder import ChartBuilder


class PowerpointReport(BaseReport):
    """Генератор отчёта в формате PowerPoint."""

    def __init__(self, template_manager: PowerpointTemplateManager, chart_builder: ChartBuilder):
        self.template_mgr = template_manager
        self.chart_builder = chart_builder
        self.renderer = PlaceholderRenderer(chart_builder)

    def generate(self, statistics: Statistics, config: ConfigModel, output_path: str) -> str:
        with self.chart_builder:
            prs = self.template_mgr.get_presentation_with_theme(config.custom_template_path)
            context = self._build_context(statistics, config)

            self._add_title_slide(prs, context)
            self._add_period_slide(prs, context, statistics)
            self._add_types_remarks_slide(prs, statistics)
            self._add_file_groups_slide(prs, statistics)      # новый слайд
            self._add_docs_with_remarks_slide(prs, statistics)
            self._add_errors_vs_a4_slide(prs, statistics)
            self._add_monthly_trend_slide(prs, statistics)
            self._add_review_leaders_slide(prs, statistics)
            self._add_errors_by_type_top10_slide(prs, statistics)
            self._add_developers_rating_slide(prs, statistics)
            self._add_conclusions1_slide(prs)
            self._add_conclusions2_slide(prs)
            self._add_specific_analysis_slide(prs, statistics)
            self._add_closing_slide(prs)

            self.renderer.render(prs, context, statistics, config)
            prs.save(output_path)
        return output_path

    def _build_context(self, stats: Statistics, config: ConfigModel) -> dict:
        total_a4 = stats.total_a4
        total_err = stats.total_errors_cat1 + stats.total_errors_cat2
        total_docs = stats.total_docs
        docs_with = stats.docs_with_errors

        if config.period_start or config.period_end:
            start_str = config.period_start.strftime('%d.%m.%Y') if config.period_start else '...'
            end_str = config.period_end.strftime('%d.%m.%Y') if config.period_end else '...'
        else:
            start_str = stats.actual_start_date.strftime('%d.%m.%Y') if stats.actual_start_date else '...'
            end_str = stats.actual_end_date.strftime('%d.%m.%Y') if stats.actual_end_date else '...'

        devs_str = ', '.join(config.selected_developers) if config.selected_developers else ''

        return {
            'total_docs': total_docs,
            'total_a4': total_a4,
            'total_remarks': total_err,
            'docs_with_remarks': docs_with,
            'critical_errors': stats.total_errors_cat1,
            'non_critical_errors': stats.total_errors_cat2,
            'critical_percent_of_a4': f'{stats.total_errors_cat1 / total_a4 * 100:.1f}%' if total_a4 else '0%',
            'non_critical_percent_of_a4': f'{stats.total_errors_cat2 / total_a4 * 100:.1f}%' if total_a4 else '0%',
            'report_date': datetime.now().strftime('%d.%m.%Y'),
            'period_start': start_str,
            'period_end': end_str,
            'selected_developers': devs_str
        }

    def _find_layout(self, prs, patterns):
        for layout in prs.slide_layouts:
            name = layout.name.lower() if layout.name else ''
            for pat in patterns:
                if pat in name:
                    return layout
        return prs.slide_layouts[0]

    def _add_title_slide(self, prs, ctx):
        layout = self._find_layout(prs, ['титул', 'title'])
        slide = prs.slides.add_slide(layout)
        if slide.shapes.title:
            slide.shapes.title.text = "Анализ результатов учета замечаний к конечному продукту ИЦ"

    def _add_period_slide(self, prs, ctx, stats):
        layout = self._find_layout(prs, ['заголовок', 'title', 'объект'])
        slide = prs.slides.add_slide(layout)
        if slide.shapes.title:
            slide.shapes.title.text = "Период и общие показатели"
        body = slide.shapes.placeholders[1] if len(slide.shapes.placeholders) > 1 else None
        if body:
            text_frame = body.text_frame
            text_frame.clear()
            lines = [
                f"Период оценки: {ctx['period_start']} – {ctx['period_end']}",
                f"Количество документов: {ctx['total_docs']}",
                f"Количество листов А4: {ctx['total_a4']}",
                f"Среднее листов на документ: {stats.total_a4 / stats.total_docs:.1f}" if stats.total_docs else ''
            ]
            if ctx.get('selected_developers'):
                lines.append(f"Отчёт по разработчикам: {ctx['selected_developers']}")
            for line in lines:
                p = text_frame.add_paragraph()
                p.text = line

    def _add_types_remarks_slide(self, prs, stats):
        layout = self._find_layout(prs, ['заголовок', 'title'])
        slide = prs.slides.add_slide(layout)
        if slide.shapes.title:
            slide.shapes.title.text = "Типы документов и распределение замечаний"
        chart_path = self.chart_builder.create_pie_chart(
            ['Категория 1', 'Категория 2'],
            [stats.total_errors_cat1, stats.total_errors_cat2],
            'Общее распределение замечаний'
        )
        slide.shapes.add_picture(chart_path, Inches(0.5), Inches(1.5), width=Inches(4))
        text_box = slide.shapes.add_textbox(Inches(5), Inches(1.5), Inches(3), Inches(5))
        tf = text_box.text_frame
        tf.word_wrap = True
        for doc_type in sorted(stats.by_type.keys()):
            p = tf.add_paragraph()
            p.text = doc_type

    def _add_file_groups_slide(self, prs, stats):
        """Слайд с гистограммой по группам файлов."""
        if not stats.docs_by_file_prefix:
            return
        layout = self._find_layout(prs, ['заголовок', 'title'])
        slide = prs.slides.add_slide(layout)
        if slide.shapes.title:
            slide.shapes.title.text = "Количество документов по группам файлов"
        groups = sorted(stats.docs_by_file_prefix.items(), key=lambda x: x[1], reverse=True)
        categories = [g[0] for g in groups]
        values = [g[1] for g in groups]
        if categories:
            chart_path = self.chart_builder.create_vertical_bar_chart(
                categories, values,
                title='',
                xlabel='Группа файлов',
                ylabel='Количество документов'
            )
            slide.shapes.add_picture(chart_path, Inches(0.5), Inches(1.5), width=Inches(8.5))

    def _add_docs_with_remarks_slide(self, prs, stats):
        layout = self._find_layout(prs, ['заголовок'])
        slide = prs.slides.add_slide(layout)
        if slide.shapes.title:
            slide.shapes.title.text = "Количество документов с замечаниями"
        docs_with = stats.docs_with_errors
        docs_without = stats.total_docs - docs_with
        chart_path = self.chart_builder.create_pie_chart(
            ['С замечаниями', 'Без замечаний'],
            [docs_with, docs_without],
            'Доля документов с замечаниями'
        )
        slide.shapes.add_picture(chart_path, Inches(1), Inches(1.5), width=Inches(5))

    def _add_errors_vs_a4_slide(self, prs, stats):
        layout = self._find_layout(prs, ['заголовок'])
        slide = prs.slides.add_slide(layout)
        if slide.shapes.title:
            slide.shapes.title.text = "Доля ошибок относительно листажа А4"
        text = (f"Всего листов А4: {stats.total_a4}\n"
                f"Ошибки кат.1: {stats.total_errors_cat1} ({self._percent(stats.total_errors_cat1, stats.total_a4)}%)\n"
                f"Ошибки кат.2: {stats.total_errors_cat2} ({self._percent(stats.total_errors_cat2, stats.total_a4)}%)")
        tb = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(4), Inches(2))
        tb.text_frame.text = text

    def _add_monthly_trend_slide(self, prs, stats):
        layout = self._find_layout(prs, ['заголовок'])
        slide = prs.slides.add_slide(layout)
        if slide.shapes.title:
            slide.shapes.title.text = "Динамика ошибок категории 1 и 2 по месяцам"
        months = sorted(stats.by_month.keys())
        if months:
            cat1 = [stats.by_month[m]['errors1'] for m in months]
            cat2 = [stats.by_month[m]['errors2'] for m in months]
            chart_path = self.chart_builder.create_monthly_trend_chart(months, cat1, cat2)
            slide.shapes.add_picture(chart_path, Inches(0.5), Inches(1.5), width=Inches(8.5))

    def _add_review_leaders_slide(self, prs, stats):
        layout = self._find_layout(prs, ['заголовок'])
        slide = prs.slides.add_slide(layout)
        if slide.shapes.title:
            slide.shapes.title.text = "Лидеры проверок документации"

    def _add_errors_by_type_top10_slide(self, prs, stats):
        layout = self._find_layout(prs, ['заголовок'])
        slide = prs.slides.add_slide(layout)
        if slide.shapes.title:
            slide.shapes.title.text = "Ошибки категории 1 по типам документов (ТОП-10)"
        types = [(t, d['errors1']) for t, d in stats.by_type.items()]
        types.sort(key=lambda x: x[1], reverse=True)
        top10 = types[:10]
        if top10:
            categories = [x[0] for x in top10]
            errors = [x[1] for x in top10]
            chart_path = self.chart_builder.create_horizontal_bar_chart(
                categories, errors, title=''
            )
            slide.shapes.add_picture(chart_path, Inches(0.5), Inches(1.5), width=Inches(8))

    def _add_developers_rating_slide(self, prs, stats):
        layout = self._find_layout(prs, ['заголовок'])
        slide = prs.slides.add_slide(layout)
        if slide.shapes.title:
            slide.shapes.title.text = "Рейтинг разработчиков по общему количеству ошибок (ТОП-10)"
        dev_list = [(dev, d['errors1'] + d['errors2']) for dev, d in stats.by_developer.items()]
        dev_list.sort(key=lambda x: x[1], reverse=True)
        top10 = dev_list[:10]
        if top10:
            rows = len(top10) + 1
            cols = 2
            left = Inches(0.5)
            top = Inches(1.5)
            width = Inches(7)
            height = Inches(0.3) * rows
            table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
            table = table_shape.table
            table.cell(0, 0).text = 'Разработчик'
            table.cell(0, 1).text = 'Всего ошибок'
            for i, (dev, total) in enumerate(top10, start=1):
                table.cell(i, 0).text = dev
                table.cell(i, 1).text = str(total)

    def _add_conclusions1_slide(self, prs):
        layout = self._find_layout(prs, ['заголовок'])
        slide = prs.slides.add_slide(layout)
        if slide.shapes.title:
            slide.shapes.title.text = "Выводы и рекомендации (часть 1)"
        tb = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5))
        tb.text_frame.text = "Содержание выводов будет добавлено позже."

    def _add_conclusions2_slide(self, prs):
        layout = self._find_layout(prs, ['заголовок'])
        slide = prs.slides.add_slide(layout)
        if slide.shapes.title:
            slide.shapes.title.text = "Выводы и рекомендации (часть 2)"
        tb = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5))
        tb.text_frame.text = "Ключевые проблемы процесса и первоочередные рекомендации."

    def _add_specific_analysis_slide(self, prs, stats):
        has_data = any(m.startswith('202') and m[5:7] in ['08', '09'] for m in stats.by_month)
        if has_data:
            layout = self._find_layout(prs, ['заголовок'])
            slide = prs.slides.add_slide(layout)
            if slide.shapes.title:
                slide.shapes.title.text = "Специфика по августу – сентябрю"

    def _add_closing_slide(self, prs):
        layout = self._find_layout(prs, ['заголовок', 'титул'])
        slide = prs.slides.add_slide(layout)
        if slide.shapes.title:
            slide.shapes.title.text = "Благодарим за внимание!"

    @staticmethod
    def _percent(part, total):
        return round(part / total * 100, 1) if total else 0.0