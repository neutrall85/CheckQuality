from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import os


class PlaceholderRenderer:
    """Заменяет плейсхолдеры вида {{...}} в презентации."""

    def __init__(self, chart_builder):
        self.chart_builder = chart_builder

    def render(self, prs: Presentation, context: dict, stats, config) -> None:
        for slide in prs.slides:
            self._process_shapes(slide.shapes, context, stats, config)

    def _process_shapes(self, shapes, context, stats, config):
        for shape in shapes:
            if shape.has_text_frame:
                self._replace_in_textframe(shape, context, stats, config)
            if shape.has_table:
                pass
            if hasattr(shape, 'shapes'):
                self._process_shapes(shape.shapes, context, stats, config)

    def _replace_in_textframe(self, shape, context, stats, config):
        text_frame = shape.text_frame
        for paragraph in text_frame.paragraphs:
            for run in paragraph.runs:
                full_text = run.text
                if not full_text:
                    continue
                if '{{' in full_text and '}}' in full_text:
                    if '{{developers_table}}' in full_text:
                        run.text = ''
                        self._insert_developers_table(text_frame, stats, shape)
                        return
                    elif '{{chart:monthly_trend}}' in full_text:
                        run.text = ''
                        self._insert_chart_monthly_trend(text_frame, stats, shape)
                        return
                    elif '{{chart:errors_by_doc_type}}' in full_text:
                        run.text = ''
                        self._insert_chart_errors_by_type(text_frame, stats, shape)
                        return
                    else:
                        for key, value in context.items():
                            placeholder = '{{' + key + '}}'
                            if placeholder in full_text:
                                full_text = full_text.replace(placeholder, str(value))
                        run.text = full_text

    def _get_slide(self, text_frame):
        """Безопасно получает объект Slide из text_frame."""
        return text_frame.shape.part.slide

    def _remove_shape(self, shape):
        """Удаляет shape со слайда."""
        try:
            shape.element.getparent().remove(shape.element)
        except Exception:
            pass

    def _insert_developers_table(self, text_frame, stats, shape):
        slide = self._get_slide(text_frame)
        dev_list = []
        for dev, data in stats.by_developer.items():
            total = data['errors1'] + data['errors2']
            dev_list.append((dev, total))
        dev_list.sort(key=lambda x: x[1], reverse=True)
        top10 = dev_list[:10]

        left = shape.left
        top = shape.top
        width = shape.width
        rows = len(top10) + 1
        height = Inches(0.3) * rows

        table_shape = slide.shapes.add_table(rows, 2, left, top, width, height)
        table = table_shape.table
        table.cell(0, 0).text = 'Разработчик'
        table.cell(0, 1).text = 'Всего ошибок (кат.1+2)'

        for i, (dev, total) in enumerate(top10, start=1):
            table.cell(i, 0).text = dev
            table.cell(i, 1).text = str(total)

        self._remove_shape(shape)

    def _insert_chart_monthly_trend(self, text_frame, stats, shape):
        slide = self._get_slide(text_frame)
        months = sorted(stats.by_month.keys())
        if months:
            cat1 = [stats.by_month[m]['errors1'] for m in months]
            cat2 = [stats.by_month[m]['errors2'] for m in months]
            chart_path = self.chart_builder.create_monthly_trend_chart(months, cat1, cat2)
            slide.shapes.add_picture(chart_path, shape.left, shape.top,
                                     width=shape.width, height=shape.height)
        self._remove_shape(shape)

    def _insert_chart_errors_by_type(self, text_frame, stats, shape):
        slide = self._get_slide(text_frame)
        types = [(t, d['errors1']) for t, d in stats.by_type.items()]
        types.sort(key=lambda x: x[1], reverse=True)
        top10 = types[:10]
        if top10:
            categories = [x[0] for x in top10]
            errors = [x[1] for x in top10]
            chart_path = self.chart_builder.create_horizontal_bar_chart(
                categories, errors, title='Ошибки категории 1 по типам документов'
            )
            slide.shapes.add_picture(chart_path, shape.left, shape.top,
                                     width=shape.width, height=shape.height)
        self._remove_shape(shape)